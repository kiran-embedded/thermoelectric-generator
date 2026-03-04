"""
PDF to PPTX Converter
Converts each page of project-report.pdf into a high-resolution
slide image inside a PowerPoint presentation.
"""
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Emu
import os, tempfile

PDF_PATH = "project-report.pdf"
PPTX_PATH = "project-report.pptx"
DPI = 300  # High resolution for crisp slides

print("📄 Converting PDF pages to high-res images...")
images = convert_from_path(PDF_PATH, dpi=DPI)
print(f"   Found {len(images)} pages.")

# A4 dimensions in EMUs (English Metric Units)
# A4 = 210mm x 297mm
SLIDE_WIDTH = Emu(7560000)   # 210mm
SLIDE_HEIGHT = Emu(10692000) # 297mm

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]

with tempfile.TemporaryDirectory() as tmpdir:
    for i, img in enumerate(images):
        print(f"   🖼️  Processing page {i+1}/{len(images)}...")
        img_path = os.path.join(tmpdir, f"page_{i+1}.png")
        img.save(img_path, "PNG")

        slide = prs.slides.add_slide(blank_layout)
        # Add image to cover full slide
        slide.shapes.add_picture(img_path, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)

print(f"💾 Saving PPTX...")
prs.save(PPTX_PATH)
print(f"✅ PPTX saved as: {PPTX_PATH}")
